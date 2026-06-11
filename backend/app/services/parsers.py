"""
Parsers that turn raw input sources (files / URLs) into plain text
or structured intermediate representations ready for chunking.

Each parser is intentionally side-effect free and raises a clear
exception on failure so the caller (sources router) can mark the
source as FAILED with a useful message.

YouTube pipeline:
  Controlled by settings.youtube_transcription_mode:
    - "whisper_audio" (default): always download the full audio with yt-dlp,
      convert to 16kHz mono WAV, chunk into ~10-minute segments, and
      transcribe each chunk locally with openai-whisper. Segment-level
      timestamps from Whisper are preserved (offset per chunk) so chunking.py
      can still produce "at MM:SS" citations.
    - "captions": fall back to the original fast path using
      youtube-transcript-api (existing captions only, with an optional
      faster-whisper fallback if captions are unavailable).
"""
from __future__ import annotations

import os
import re
import tempfile
import uuid
from typing import Optional

import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader
from pptx import Presentation

from app.config import get_settings

_settings = get_settings()

# youtube-transcript-api 
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api._errors import (
    CouldNotRetrieveTranscript,
    IpBlocked,
    NoTranscriptFound,
    RequestBlocked,
    TranscriptsDisabled,
    VideoUnavailable,
    YouTubeDataUnparsable,
    YouTubeRequestFailed,
    YouTubeTranscriptApiException,
)


try:
    from youtube_transcript_api._errors import VideoUnplayable
    _HAS_VIDEO_UNPLAYABLE = True
except ImportError:
    VideoUnplayable = None 
    _HAS_VIDEO_UNPLAYABLE = False

REQUEST_TIMEOUT = 20

_BROWSER_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/124.0.0.0 Safari/537.36"
    ),
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.9",
    "Accept-Encoding": "gzip, deflate, br",
    "Connection": "keep-alive",
    "Upgrade-Insecure-Requests": "1",
}

_OEMBED_UA = (
    "Mozilla/5.0 (compatible; SamasocialLearningAssistant/1.0; "
    "+https://samasocial.in)"
)



# pdf
def parse_pdf(file_path: str) -> dict:
    """
    Returns {"pages": [{"page_number": int, "text": str}], "title": str}
    """
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append({"page_number": i, "text": text})

    if not pages:
        raise ValueError(
            "No extractable text found in PDF "
            "(it may be a scanned/image-only document)."
        )

    title = ""
    if reader.metadata and reader.metadata.title:
        title = reader.metadata.title.strip()

    return {"pages": pages, "title": title}


# PPTX
def parse_pptx(file_path: str) -> dict:
    """
    Returns {"slides": [{"slide_number": int, "text": str}], "title": str}
    """
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                texts.append(f"[Speaker notes] {note}")

        slide_text = "\n".join(texts).strip()
        if slide_text:
            slides.append({"slide_number": i, "text": slide_text})

    if not slides:
        raise ValueError("No extractable text found in the PPTX file.")

    return {"slides": slides, "title": ""}


# YouTube helpers
_YOUTUBE_ID_PATTERNS = [
    r"(?:v=|\/)([0-9A-Za-z_-]{11}).*",
    r"youtu\.be\/([0-9A-Za-z_-]{11})",
]


def extract_youtube_id(url: str) -> str:
    for pattern in _YOUTUBE_ID_PATTERNS:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    raise ValueError("Could not extract a video ID from the given YouTube URL.")


def _fetch_youtube_title(video_id: str) -> str:
    """Best-effort title fetch via oEmbed — returns '' on any failure."""
    try:
        resp = requests.get(
            "https://www.youtube.com/oembed",
            params={
                "url": f"https://www.youtube.com/watch?v={video_id}",
                "format": "json",
            },
            headers={"User-Agent": _OEMBED_UA},
            timeout=REQUEST_TIMEOUT,
        )
        if resp.ok:
            return resp.json().get("title", "")
    except requests.RequestException:
        pass
    return ""


def _captions_to_segments(fetched) -> list[dict]:
    """
    Convert a FetchedTranscript (youtube-transcript-api v1.x) to the
    canonical segment format: [{"start": float, "duration": float, "text": str}]
    """
    return fetched.to_raw_data()


# ─────────────────────────────────────────────────────────────────────────── #
# Full-audio Whisper pipeline (default YouTube path)
# ─────────────────────────────────────────────────────────────────────────── #

_whisper_model_cache = {}


def _get_whisper_model():
    """Load (and cache) the local openai-whisper model named by settings.whisper_model."""
    try:
        import whisper  # noqa: PLC0415
    except ImportError as exc:
        raise ValueError(
            "openai-whisper is not installed. Run: pip install openai-whisper "
            "and ensure ffmpeg is on your PATH."
        ) from exc

    model_size = getattr(_settings, "whisper_model", "small")
    if model_size not in _whisper_model_cache:
        _whisper_model_cache[model_size] = whisper.load_model(model_size)
    return _whisper_model_cache[model_size]


def _download_youtube_audio(video_id: str, tmp_dir: str) -> str:
    """
    Download the best-audio stream for a YouTube video into tmp_dir using yt-dlp.
    Returns the path to the downloaded file.

    Mirrors the AVA repo's approach: tries progressively more permissive format
    selectors so we never hit "format not available", then scans the temp dir
    for whatever yt-dlp actually wrote.
    """
    try:
        import yt_dlp  # noqa: PLC0415
    except ImportError as exc:
        raise ValueError(
            "yt-dlp is not installed. Run: pip install yt-dlp"
        ) from exc

    url = f"https://www.youtube.com/watch?v={video_id}"
    output_template = os.path.join(tmp_dir, "%(id)s.%(ext)s")

    # Tiered format selectors — tried in order until one succeeds.
    # DO NOT add extractor_args skip:dash/hls — that strips most available
    # formats and causes "Requested format is not available" errors.
    FORMAT_TIERS = [
        # Tier 1: audio-only streams (smallest, fastest download)
        "bestaudio[ext=m4a]/bestaudio[ext=webm]/bestaudio",
        # Tier 2: any format that has an audio codec
        "best[acodec!=none]",
        # Tier 3: absolute fallback — whatever yt-dlp considers "best"
        "best",
    ]

    base_opts = {
        "outtmpl": output_template,
        "quiet": True,
        "no_warnings": True,
        "overwrites": True,
        "nopart": True,
    }

    last_exc: Exception | None = None
    for fmt in FORMAT_TIERS:
        ydl_opts = {**base_opts, "format": fmt}
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                filename = ydl.prepare_filename(info)

            # yt-dlp sometimes changes the extension after post-processing;
            # scan tmp_dir for the actual file if prepare_filename is stale.
            if not filename or not os.path.exists(filename):
                candidates = sorted([
                    os.path.join(tmp_dir, f)
                    for f in os.listdir(tmp_dir)
                    if f.startswith(info.get("id", ""))
                ], key=os.path.getsize, reverse=True)
                if candidates:
                    filename = candidates[0]

            if filename and os.path.exists(filename):
                return filename

        except Exception as exc:  # noqa: BLE001
            last_exc = exc
            continue  # try next tier

    # Final fallback: let yt-dlp choose automatically.
    try:
        with yt_dlp.YoutubeDL(base_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            filename = ydl.prepare_filename(info)
            if filename and os.path.exists(filename):
                return filename
    except Exception as exc:
        last_exc = exc

    raise ValueError(
        f"yt-dlp could not download audio for video '{video_id}'. "
        f"Last error: {last_exc}. "
        "Make sure yt-dlp is up-to-date: pip install -U yt-dlp"
    )


def _convert_to_wav(input_path: str) -> str:
    """Convert any audio/video file to 16kHz mono WAV using pydub."""
    try:
        from pydub import AudioSegment  # noqa: PLC0415
    except ImportError as exc:
        raise ValueError(
            "pydub is not installed. Run: pip install pydub "
            "and ensure ffmpeg is on your PATH."
        ) from exc

    output_path = os.path.splitext(input_path)[0] + "_converted.wav"
    audio = AudioSegment.from_file(input_path)
    audio = audio.set_channels(1).set_frame_rate(16000)
    audio.export(output_path, format="wav")
    return output_path


def _chunk_audio(wav_path: str, chunk_minutes: int) -> list[str]:
    """Split a WAV file into ~chunk_minutes-long pieces, returning their paths."""
    from pydub import AudioSegment  # noqa: PLC0415

    audio = AudioSegment.from_wav(wav_path)
    chunk_ms = chunk_minutes * 60 * 1000

    chunk_paths = []
    for i, start in enumerate(range(0, len(audio), chunk_ms)):
        chunk = audio[start: start + chunk_ms]
        chunk_path = f"{wav_path}_chunk_{i}.wav"
        chunk.export(chunk_path, format="wav")
        chunk_paths.append(chunk_path)

    return chunk_paths


def _transcribe_chunk_whisper(chunk_path: str) -> list[dict]:
    """
    Transcribe a single audio chunk with openai-whisper.
    Returns the raw list of Whisper segment dicts (each has start/end/text).
    """
    model = _get_whisper_model()
    result = model.transcribe(chunk_path, task="transcribe")
    return result.get("segments", [])


def _whisper_full_audio_transcribe(video_id: str) -> list[dict]:
    """
    Full pipeline: download -> convert to WAV -> chunk -> transcribe each
    chunk with Whisper -> merge segments with corrected (offset) timestamps.

    Returns segments in {"start": float, "duration": float, "text": str} format,
    with "start" expressed relative to the full video.

    Raises ValueError with a human-readable message on any failure.
    """
    tmp_dir = tempfile.mkdtemp(prefix="yt_audio_")
    downloaded_path: Optional[str] = None
    wav_path: Optional[str] = None
    chunk_paths: list[str] = []

    try:
        downloaded_path = _download_youtube_audio(video_id, tmp_dir)
        wav_path = _convert_to_wav(downloaded_path)

        chunk_minutes = getattr(_settings, "whisper_chunk_minutes", 10)
        chunk_paths = _chunk_audio(wav_path, chunk_minutes)

        if not chunk_paths:
            raise ValueError("Audio chunking produced no chunks for this video.")

        segments: list[dict] = []
        chunk_offset_seconds = 0.0

        for chunk_path in chunk_paths:
            chunk_segments = _transcribe_chunk_whisper(chunk_path)

            for seg in chunk_segments:
                start = float(seg["start"]) + chunk_offset_seconds
                end = float(seg["end"]) + chunk_offset_seconds
                text = seg.get("text", "").strip()
                if not text:
                    continue
                segments.append(
                    {
                        "start": start,
                        "duration": max(end - start, 0.0),
                        "text": text,
                    }
                )

            chunk_offset_seconds += chunk_minutes * 60

        if not segments:
            raise ValueError("Whisper transcription produced no output for this video.")

        return segments

    except ValueError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"Whisper full-audio transcription failed: {exc}") from exc
    finally:
        # Clean up all temp files (downloaded audio, wav, chunk wavs)
        for path in [downloaded_path, wav_path, *chunk_paths]:
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except OSError:
                    pass
        try:
            os.rmdir(tmp_dir)
        except OSError:
            pass  # dir not empty (or already gone) is fine


# ─────────────────────────────────────────────────────────────────────────── #
# Sarvam AI transcription pipeline  (Hindi / Hinglish → English)
# ─────────────────────────────────────────────────────────────────────────── #

# Sarvam's sync STT-translate API rejects audio longer than 30s.
# We slice each chunk into 25s pieces (with a 5s safety margin).
_SARVAM_PIECE_SECONDS = 25
_SARVAM_STT_TRANSLATE_URL = "https://api.sarvam.ai/speech-to-text-translate"


def _send_piece_to_sarvam(piece_path: str, api_key: str, model: str) -> str:
    """POST one ≤30s WAV piece to Sarvam and return the English transcript."""
    headers = {"api-subscription-key": api_key}
    with open(piece_path, "rb") as f:
        files = {"file": (os.path.basename(piece_path), f, "audio/wav")}
        data = {"model": model, "with_diarization": "false"}
        response = requests.post(
            _SARVAM_STT_TRANSLATE_URL,
            headers=headers,
            files=files,
            data=data,
            timeout=120,
        )

    if not response.ok:
        raise ValueError(
            f"Sarvam API returned {response.status_code}: {response.text}"
        )

    return response.json().get("transcript", "")


def _transcribe_chunk_sarvam(chunk_path: str, api_key: str, model: str) -> str:
    """
    Sarvam sync API only accepts ≤30s audio. We split this 10-min chunk into
    25-second pieces, send each separately, and join the transcripts.
    Returns plain text (Sarvam translates Hindi → English inline).
    """
    from pydub import AudioSegment  # noqa: PLC0415

    audio = AudioSegment.from_wav(chunk_path)
    piece_ms = _SARVAM_PIECE_SECONDS * 1000
    total_pieces = (len(audio) + piece_ms - 1) // piece_ms

    full_text = ""
    for i, start in enumerate(range(0, len(audio), piece_ms)):
        piece = audio[start: start + piece_ms]
        piece_path = f"{chunk_path}_sv_{i}.wav"
        piece.export(piece_path, format="wav")
        try:
            print(f"  → Sarvam piece {i + 1}/{total_pieces} ...")
            full_text += _send_piece_to_sarvam(piece_path, api_key, model) + " "
        finally:
            if os.path.exists(piece_path):
                try:
                    os.remove(piece_path)
                except OSError:
                    pass

    return full_text.strip()


def _sarvam_full_audio_transcribe(video_id: str) -> list[dict]:
    """
    Full pipeline for Hindi/Hinglish videos:
      download → convert to 16kHz mono WAV → chunk (~10 min) →
      send each chunk to Sarvam AI in 25s pieces → reassemble with timestamps.

    Returns segments in {"start": float, "duration": float, "text": str} format.
    Raises ValueError with a human-readable message on any failure.
    """
    api_key = getattr(_settings, "sarvam_api_key", "") or os.getenv("SARVAM_API_KEY", "")
    if not api_key:
        raise ValueError(
            "SARVAM_API_KEY is not set. Add it to your .env file to use "
            "Hindi/Hinglish transcription."
        )

    sarvam_model = getattr(_settings, "sarvam_stt_model", "saaras:v2.5")
    chunk_minutes = getattr(_settings, "whisper_chunk_minutes", 10)

    tmp_dir = tempfile.mkdtemp(prefix="yt_sarvam_")
    downloaded_path: Optional[str] = None
    wav_path: Optional[str] = None
    chunk_paths: list[str] = []

    try:
        downloaded_path = _download_youtube_audio(video_id, tmp_dir)
        wav_path = _convert_to_wav(downloaded_path)
        chunk_paths = _chunk_audio(wav_path, chunk_minutes)

        if not chunk_paths:
            raise ValueError("Audio chunking produced no chunks for this video.")

        segments: list[dict] = []
        chunk_offset_seconds = 0.0

        for idx, chunk_path in enumerate(chunk_paths):
            print(f"  Sarvam: transcribing chunk {idx + 1}/{len(chunk_paths)} ...")
            text = _transcribe_chunk_sarvam(chunk_path, api_key, sarvam_model)

            if text:
                segments.append(
                    {
                        "start": chunk_offset_seconds,
                        "duration": float(chunk_minutes * 60),
                        "text": text,
                    }
                )

            chunk_offset_seconds += chunk_minutes * 60

        if not segments:
            raise ValueError("Sarvam transcription produced no output for this video.")

        return segments

    except ValueError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"Sarvam transcription failed: {exc}") from exc
    finally:
        for path in [downloaded_path, wav_path, *chunk_paths]:
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except OSError:
                    pass
        try:
            os.rmdir(tmp_dir)
        except OSError:
            pass



# Whisper fallback
def _whisper_transcribe(video_id: str) -> list[dict]:
    """
    Download audio with yt-dlp and transcribe with faster-whisper.
    Returns segments in {"start": float, "duration": float, "text": str} format.

    Raises ValueError with a human-readable message on any failure.
    """
    try:
        import yt_dlp  # noqa: PLC0415
    except ImportError as exc:
        raise ValueError(
            "yt-dlp is not installed. Run: pip install yt-dlp"
        ) from exc

    try:
        from faster_whisper import WhisperModel  # noqa: PLC0415
    except ImportError as exc:
        raise ValueError(
            "faster-whisper is not installed. Run: pip install faster-whisper  "
            "and ensure ffmpeg is on your PATH."
        ) from exc

    tmp_dir = tempfile.mkdtemp(prefix="yt_audio_")
    audio_path: Optional[str] = None

    try:
        url = f"https://www.youtube.com/watch?v={video_id}"
        output_template = os.path.join(tmp_dir, "%(id)s.%(ext)s")

        # Same tiered format approach as _download_youtube_audio —
        # avoids "Requested format is not available" on DASH-only videos.
        FORMAT_TIERS = [
            "bestaudio[ext=m4a]/bestaudio[ext=webm]/bestaudio",
            "best[acodec!=none]",
            "best",
        ]

        base_opts = {
            "outtmpl": output_template,
            "quiet": True,
            "no_warnings": True,
            "overwrites": True,
            "nopart": True,
            "postprocessors": [
                {
                    "key": "FFmpegExtractAudio",
                    "preferredcodec": "mp3",
                    "preferredquality": "128",
                }
            ],
        }

        info = None
        last_dl_exc: Exception | None = None
        for fmt in FORMAT_TIERS:
            try:
                with yt_dlp.YoutubeDL({**base_opts, "format": fmt}) as ydl:
                    info = ydl.extract_info(url, download=True)
                break
            except Exception as exc:  # noqa: BLE001
                last_dl_exc = exc
                continue

        if info is None:
            try:
                with yt_dlp.YoutubeDL(base_opts) as ydl:
                    info = ydl.extract_info(url, download=True)
            except Exception as exc:
                last_dl_exc = exc

        if info is None:
            raise ValueError(
                f"yt-dlp could not download audio. Last error: {last_dl_exc}. "
                "Try: pip install -U yt-dlp"
            )

        # After FFmpeg post-processing the extension is always mp3
        audio_path = os.path.join(tmp_dir, f"{info['id']}.mp3")

        if not audio_path or not os.path.exists(audio_path):
            # Fallback: find any audio file that was created
            for fname in os.listdir(tmp_dir):
                if fname.startswith(info["id"]):
                    audio_path = os.path.join(tmp_dir, fname)
                    break

        if not audio_path or not os.path.exists(audio_path):
            raise ValueError("yt-dlp downloaded the audio but the file could not be located.")

        # Load Whisper model (downloaded once and cached by faster-whisper)
        model_size = getattr(_settings, "whisper_model_size", "base")
        model = WhisperModel(model_size, device="cpu", compute_type="int8")

        whisper_segments, _ = model.transcribe(
            audio_path,
            beam_size=5,
            word_timestamps=False,
            vad_filter=True,         
        )

        segments: list[dict] = []
        for seg in whisper_segments:
            duration = seg.end - seg.start
            segments.append(
                {
                    "start": float(seg.start),
                    "duration": float(duration),
                    "text": seg.text.strip(),
                }
            )

        if not segments:
            raise ValueError("Whisper transcription produced no output for this video.")

        return segments

    except ValueError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"Whisper transcription failed: {exc}") from exc
    finally:
        # Always clean up temp audio files
        if audio_path and os.path.exists(audio_path):
            try:
                os.remove(audio_path)
            except OSError:
                pass
        try:
            os.rmdir(tmp_dir)
        except OSError:
            pass  # dir not empty is fine; we only committed to removing the audio


def parse_youtube(url: str, language: str = "english") -> dict:
    """
    Returns:
        {
            "segments":         [{"start": float, "duration": float, "text": str}],
            "title":            str,
            "video_id":         str,
            "processing_method":"captions" | "whisper" | "whisper_full_audio" | "sarvam",
        }

    Strategy depends on `language` param and settings.youtube_transcription_mode:
      - language="hinglish" (or "hindi"): use Sarvam AI for STT + English
        translation. Downloads full audio with yt-dlp, chunks into ~10-min
        pieces, sends each in 25s slices to the Sarvam API.
      - language="english" + mode="whisper_audio" (default): download full
        audio with yt-dlp, convert to 16kHz mono WAV, chunk, and transcribe
        locally with openai-whisper. Works for any video regardless of
        caption availability.
      - language="english" + mode="captions": use youtube-transcript-api
        (instant), falling back to faster-whisper if captions are unavailable
        and youtube_whisper_fallback_enabled is True.
    """
    video_id = extract_youtube_id(url)

    # ── Sarvam path: Hindi / Hinglish videos ────────────────────────────────
    if language.lower() in ("hinglish", "hindi"):
        print("Using Sarvam AI for Hindi/Hinglish transcription...")
        segments = _sarvam_full_audio_transcribe(video_id)
        title = _fetch_youtube_title(video_id)
        return {
            "segments": segments,
            "title": title,
            "video_id": video_id,
            "processing_method": "sarvam",
        }

    # ── Default / preferred path: full-audio Whisper transcription ─────────
    if _settings.youtube_transcription_mode == "whisper_audio":
        segments = _whisper_full_audio_transcribe(video_id)
        title = _fetch_youtube_title(video_id)
        return {
            "segments": segments,
            "title": title,
            "video_id": video_id,
            "processing_method": "whisper_full_audio",
        }

    # ── Optional fast path: captions, with faster-whisper fallback 
    caption_error: Optional[str] = None

    try:
        api = YouTubeTranscriptApi()
        transcript_list = api.list(video_id)

        try:
            transcript = transcript_list.find_transcript(["en", "en-US", "en-GB"])
        except NoTranscriptFound:
            # Fall back to first available language, translated if possible
            transcript = next(iter(transcript_list))
            if transcript.is_translatable:
                try:
                    transcript = transcript.translate("en")
                except Exception:
                    pass

        fetched = transcript.fetch()
        segments = _captions_to_segments(fetched)

        if not segments:
            raise ValueError("Caption transcript was empty.")

        title = _fetch_youtube_title(video_id)
        return {
            "segments": segments,
            "title": title,
            "video_id": video_id,
            "processing_method": "captions",
        }

    # Collect the caption failure reason — we may still recover via Whisper
    except TranscriptsDisabled:
        caption_error = "Transcripts/captions are disabled for this video."
    except VideoUnavailable:
        caption_error = "Video is unavailable (private, deleted, or region-locked)."
    except NoTranscriptFound:
        caption_error = "No captions are available for this video."
    except (RequestBlocked, IpBlocked):
        caption_error = (
            "YouTube is blocking caption requests from this server/IP. "
            "Falling back to audio transcription."
        )
    except YouTubeRequestFailed as exc:
        caption_error = f"YouTube caption request failed: {exc}."
    except YouTubeDataUnparsable as exc:
        caption_error = f"YouTube returned unparsable caption data: {exc}."
    except CouldNotRetrieveTranscript:
        caption_error = (
            "Could not retrieve captions — YouTube may be rate-limiting this server."
        )
    except YouTubeTranscriptApiException as exc:
        caption_error = f"Caption API error: {exc}."
    except ValueError as exc:
        caption_error = str(exc)
    except Exception as exc:  # noqa: BLE001
        caption_error = f"Unexpected caption error: {exc}."

    # ── Step 2: faster-whisper fallback
    if not _settings.youtube_whisper_fallback_enabled:
        raise ValueError(
            f"Caption extraction failed ({caption_error}) and the Whisper "
            "fallback is disabled. Enable it by setting "
            "YOUTUBE_WHISPER_FALLBACK_ENABLED=true in your .env, then ensure "
            "ffmpeg and faster-whisper are installed."
        )

    try:
        segments = _whisper_transcribe(video_id)
    except ValueError as exc:
        raise ValueError(
            f"Caption extraction failed: {caption_error}  |  "
            f"Whisper fallback also failed: {exc}"
        ) from exc

    title = _fetch_youtube_title(video_id)
    return {
        "segments": segments,
        "title": title,
        "video_id": video_id,
        "processing_method": "whisper",
    }


# Webpage
_NOISE_TAGS = [
    "script", "style", "nav", "header", "footer",
    "aside", "form", "noscript", "svg", "iframe",
]


def parse_webpage(url: str) -> dict:
    """Returns {"text": str, "title": str}"""
    try:
        resp = requests.get(url, headers=_BROWSER_HEADERS, timeout=REQUEST_TIMEOUT)
    except requests.exceptions.ConnectionError as exc:
        raise ValueError(
            f"Could not connect to '{url}'. Check the URL or try again later."
        ) from exc
    except requests.exceptions.Timeout as exc:
        raise ValueError(
            f"Request to '{url}' timed out after {REQUEST_TIMEOUT}s."
        ) from exc
    except requests.exceptions.RequestException as exc:
        raise ValueError(f"Failed to fetch webpage: {exc}") from exc

    if resp.status_code == 403:
        raise ValueError(
            f"Access denied (HTTP 403) for '{url}'. "
            "The site blocks automated requests — try a different URL."
        )
    if resp.status_code == 404:
        raise ValueError(f"Page not found (HTTP 404): '{url}'.")
    if not resp.ok:
        raise ValueError(f"HTTP {resp.status_code} error fetching '{url}'.")

    soup = BeautifulSoup(resp.text, "html.parser")

    for tag in soup(_NOISE_TAGS):
        tag.decompose()

    title = ""
    if soup.title and soup.title.string:
        title = soup.title.string.strip()
    if not title:
        title = url

    main_content = (
        soup.find("article")
        or soup.find("main")
        or soup.find(id=re.compile(r"(content|main|article)", re.I))
        or soup.find(class_=re.compile(r"(content|main|article|post|entry)", re.I))
        or soup.body
        or soup
    )

    text = main_content.get_text(separator="\n")
    lines = [line.strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    cleaned = "\n\n".join(lines)

    if len(cleaned) < 50:
        raise ValueError(
            f"Could not extract readable text from '{url}'. "
            "The page may require JavaScript to render, or may be blocking "
            "automated access. Try a different URL."
        )

    return {"text": cleaned, "title": title}